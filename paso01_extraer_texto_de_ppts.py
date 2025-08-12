from pptx import Presentation
import yaml
import subprocess
import os
from PIL import Image
import time
import pandas as pd
import numpy  as np

ti = time.time()

path_entrada = f'data/PPTs_originales'
path_salida = f'data/PPTs_elementos'
path_slides = f'data/PPTs_separados'
nombres_archivos = os.popen(f'ls {path_entrada}').read().split('\n')[:-1]
#Para eliminar archivos temporales
nombres_archivos = [nombre for nombre in nombres_archivos if nombre[0]!='~'] 

# Cargar la configuración desde el archivo YAML
with open(f'config.yml', 'r') as file:
    config = yaml.safe_load(file)

db_name = config['database']['name']
unidad = config['extrae_ppt']['unidad']
conversiones = {'centímetros': 360000, 'pulgadas': 914400, 'pixeles': 9525}
conversion = conversiones[unidad]

# Diccionario de colores comunes en español y sus valores RGB
colores_rgb = {
    "rojo": (255, 0, 0),
    "azul": (0, 0, 255),
    "verde": (0, 255, 0),
    "amarillo": (255, 255, 0),
    "negro": (0, 0, 0),
    "blanco": (255, 255, 255),
    "gris": (128, 128, 128),
    "morado": (128, 0, 128),
    "naranja": (255, 165, 0),
    "rosado": (255, 192, 203),
    "marrón": (165, 42, 42),
    "cian": (0, 255, 255),
    "magenta": (255, 0, 255),
    "verde oscuro": (0, 100, 0),
    "azul claro": (173, 216, 230),
    "verde lima": (50, 205, 50),
    "celeste": (135, 206, 235),
    "turquesa": (64, 224, 208),
    "violeta": (238, 130, 238),
    "dorado": (255, 215, 0),
    "": (-255, -255, -255)
}

def color_mas_cercano(rgb):
    min_distance = float('inf')
    closest_color = None
    for color_name, color_rgb in colores_rgb.items():
        distance = (color_rgb[0] - rgb[0]) ** 2 + (color_rgb[1] - rgb[1]) ** 2 + (color_rgb[2] - rgb[2]) ** 2
        if distance < min_distance:
            min_distance = distance
            closest_color = color_name
    return closest_color

def run_bash_command(command):
    try:
        subprocess.run(command, shell=True, capture_output=True, text=True, check=True)
    except subprocess.CalledProcessError as e:
        print(f"Error Obtenido: {e.stderr}")
    except Exception as e:
        print("Ocurrió un error inesperado:")
        print(str(e))

def muestra_atributos(objeto):
    for atributo in dir(objeto):
        # Filtrar solo los atributos definidos por el usuario, excluyendo métodos internos
        if not atributo.startswith('__'):
            try:
                valor = getattr(objeto, atributo)
                # Evitar métodos y funciones llamables
                if not callable(valor):
                    print(f'Atributo: {atributo}, Valor: {valor}')
            except Exception as e:
                # Capturar y mostrar excepciones sin detener el loop
                print(f'Error al obtener {atributo}: {e}')

def imprime_texto(lista_textos):
    for texto in lista_textos:
        print(texto)

def obtener_posicion_relativa(rango_x, rango_y, ancho_slide, alto_slide):
    """
    Determina la posición relativa de un elemento dentro de una diapositiva, considerando su
    posición central y la extensión de sus rangos en los ejes.

    :param rango_x: Tupla con los valores mínimo y máximo de x (min_x, max_x).
    :param rango_y: Tupla con los valores mínimo y máximo de y (min_y, max_y).
    :param ancho_slide: Ancho de la diapositiva.
    :param alto_slide: Alto de la diapositiva.
    :return: String con la posición relativa del elemento.
    """

    # Definir divisiones en tercios
    tercio_x = ancho_slide / 3
    tercio_y = alto_slide / 3

    # Evaluar la cobertura horizontal
    if rango_x[0] <= tercio_x and rango_x[1] >= 2 * tercio_x:
        pos_x = "completo"
    elif rango_x[0] <= tercio_x:
        pos_x = "izquierda"
    elif rango_x[1] >= 2 * tercio_x:
        pos_x = "derecha"
    else:
        pos_x = "centro"

    # Evaluar la cobertura vertical
    if rango_y[0] <= tercio_y and rango_y[1] >= 2 * tercio_y:
        pos_y = "completo"
    elif rango_y[0] <= tercio_y:
        pos_y = "arriba"
    elif rango_y[1] >= 2 * tercio_y:
        pos_y = "abajo"
    else:
        pos_y = "centro"

    # Combinar la posición vertical y horizontal, y manejar casos de coberturas completas
    if pos_x == "completo" and pos_y == "completo":
        posicion_relativa = "cobertura completa"
    elif pos_x == "completo":
        posicion_relativa = f"{pos_y} completo"
    elif pos_y == "completo":
        posicion_relativa = f"completo {pos_x}"
    else:
        posicion_relativa = f"{pos_y}-{pos_x}"

    return posicion_relativa


def guarda_wmf_shape_en_png(wmf_shape, output_path):
    temp_wmf_path = output_path.replace('.png', '.wmf')
    with open(temp_wmf_path, 'wb') as f:
        f.write(wmf_shape.image.blob)

def convertir_lote_wmf_a_png(wmf_files):
    comando = ['soffice', '--headless', '--convert-to', 'png', '--outdir',
               os.path.dirname(wmf_files[0])] + wmf_files

    try:
        subprocess.run(comando, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        print("Conversión completada para el lote.")
    except subprocess.CalledProcessError as e:
        print(f"Error durante la conversión en lote: {e}")

def corregir_orientacion_imagen(path):
    # Abrir la imagen, invertir horizontalmente y guardar de nuevo
    with Image.open(path) as img:
        img_invertida = img.transpose(Image.FLIP_LEFT_RIGHT)
        img_invertida.save(path)

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slides_text = {}

    n_slide=0
    for slide in prs.slides:
        n_slide += 1
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_text[n_slide] = slide_text

    return slides_text

def devuelve_atributo_presente_o_ausente(objeto, atributos, valor_por_defecto):
    """
    Intenta obtener un valor de un objeto basado en una cadena de atributos.
    Si ocurre cualquier excepción, retorna un valor por defecto.
    :param objeto: El objeto del cual se quiere obtener el valor.
    :param atributos: Cadena de atributos a obtener, separados por puntos (e.g., 'fill.fore_color.rgb').
    :param valor_por_defecto: El valor que se retorna si ocurre una excepción.
    :return: El valor obtenido del objeto o el valor por defecto si ocurre una excepción."""

    try:
        # Divide la cadena de atributos y accede a ellos secuencialmente
        for attr in atributos.split('.'):
            objeto = getattr(objeto, attr)
        return objeto
    except Exception:
        return valor_por_defecto


def describe_un_color(color_relleno, color_borde):
    if color_relleno=='':
        st_relleno = 'sin color'
    else:
        st_relleno = f'color {color_relleno}'

    if color_borde=='':
        st_borde = 'sin color'
    else:
        st_borde = f'color {color_borde}'
    
    if color_relleno == color_borde:
        st_final = f'de relleno y borde {st_relleno}'
    else:
        st_final = f'de relleno {st_relleno} y borde {st_borde}'
    return st_final


def describe_colores(diccionario_colores):
    # Esta función recibe un diccionario con clave la id de los textos
    # y valor una tupla de los colores (color relleno, color borde).
    # Y devuelve primero la descripción general de los colores del grupo más
    # un diccionario con los valores excepcionales si aplica
    lista_tupla_colores = list(diccionario_colores.values())
    if len(diccionario_colores)==1:
        general=describe_un_color(*lista_tupla_colores[0])
        especificos = {}
    
    elif len(diccionario_colores)==2:
        if lista_tupla_colores[0]==lista_tupla_colores[1]:
            general=describe_un_color(*lista_tupla_colores[0])
            especificos = {}
        else:
            general = 'distintos colores'
            especificos = {key: describe_un_color(*el) for key, el in diccionario_colores.items()}
    
    else:
        colores_diferentes = set(lista_tupla_colores)
        ncolores_diferentes = len(colores_diferentes)
        if ncolores_diferentes == 1:
            general = describe_un_color(*lista_tupla_colores[0])
            especificos = {}
        elif ncolores_diferentes == 2: # Es decir hay un color excepcional
            color_repetido = [color for color in lista_tupla_colores if lista_tupla_colores.count(color)>1][0]
            general = f'{describe_un_color(*color_repetido)} con excepcion de 1 texto'
            especificos = {key: describe_un_color(*color) for key, color in diccionario_colores.items()
                           if color!=color_repetido}
        else:
            general = 'diferentes colores'
            especificos = {}

    return general, especificos

def traslapes_imagen_texto(shapes_texto, shapes_imagen):
    datos_traslape = []
    for nombre_imagen, shape_imagen in shapes_imagen.items():
        rango_x_img, rango_y_img = shape_imagen['Rango X'], shape_imagen['Rango Y']
        for nombre_texto, shape_texto in shapes_texto.items():
            rango_x_text, rango_y_text = shape_texto['Rango X'], shape_texto['Rango Y']
            if (rango_x_img[0] <= rango_x_text[0] <= rango_x_img[1] and 
                rango_x_img[0] <= rango_x_text[1] <= rango_x_img[1] and
                rango_y_img[0] <= rango_y_text[0] <= rango_y_img[1] and
                rango_y_img[0] <= rango_y_text[1] <= rango_y_img[1]):
                orden_texto, orden_imagen = shape_texto['Orden'], shape_imagen['Orden']
                encima = 'texto' if orden_texto > orden_imagen else 'imagen'
                datos_traslape.append([nombre_texto, nombre_imagen, encima])
    traslapes = pd.DataFrame(datos_traslape, columns=['Texto', 'Imagen', 'Encima'])

    imagenes_traslapantes = sorted(list(set(traslapes['Imagen'])))
    traslapes_mayoritarios = dict([])
    for imagen in imagenes_traslapantes:
        n_traslapes = len(traslapes['Imagen']==imagen)
        if n_traslapes > len(shapes_texto)/2: # Esto significa que traslapa con la mayoria de los textos
            tipos_traslape=traslapes[traslapes['Imagen']==imagen]['Encima'].values
            if len(np.where(tipos_traslape=='texto')[0]) >= len(shapes_texto)/2:
                tipo_traslape = 'debajo'
            elif len(np.where(tipos_traslape=='imagen')[0]) >= len(shapes_texto)/2:
                tipo_traslape = 'sobre'
            else:
                tipo_traslape = 'sobre y debajo'
            traslapes_mayoritarios[imagen] = tipo_traslape


    return traslapes, traslapes_mayoritarios



def describe_text_shape(shape, n_shape):
    dict = {
        "Texto": shape.text,
        "Rango X": (shape.left/conversion,shape.left/conversion + shape.width/conversion),
        "Rango Y": (shape.top/conversion, shape.top/conversion + shape.height/conversion),
        "Nombre": shape.name,
        "Rotación": shape.rotation,
        "Color Relleno": color_mas_cercano(devuelve_atributo_presente_o_ausente(shape, 'fill.fore_color.rgb', (-255, -255, -255))),
        "Color Borde": color_mas_cercano(devuelve_atributo_presente_o_ausente(shape, 'line.color.rgb', (-255, -255, -255))),
        "Orden": n_shape
    }
    return dict

def agrupar_text_shapes(text_shapes, ancho_slide, alto_slide, porcentaje_umbral=1):
    grupos = {}
    grupo_id = 1
    umbral_x = ancho_slide * porcentaje_umbral/100  # Umbral de distancia minima de grupos

    # Ordenar los elementos por su valor mínimo de x
    elementos_ordenados = sorted(text_shapes.items(), key=lambda item: item[1]['Rango X'][0])

    # Inicializar el primer grupo
    grupo_actual = [elementos_ordenados[0]]

    # Agrupar por la distancia al vecino en x
    for i in range(1, len(elementos_ordenados)):
        actual = elementos_ordenados[i]
        anterior = elementos_ordenados[i - 1]

        # Comparar distancia del mínimo de x con el anterior
        if abs(actual[1]['Rango X'][0] - anterior[1]['Rango X'][0]) < umbral_x:
            grupo_actual.append(actual)
        else:
            # Crear un nuevo grupo y reiniciar
            grupos[f'Grupo{grupo_id}'] = grupo_actual
            grupo_actual = [actual]
            grupo_id += 1

    # Añadir el último grupo
    grupos[f'Grupo{grupo_id}'] = grupo_actual

    # Reagrupar los elementos que quedaron solos por alineación central usando el mismo método
    elementos_solitarios = {k: v for k, v in grupos.items() if len(v) == 1}
    for grupo, elementos in elementos_solitarios.items():
        if len(elementos) == 1:
            # Ordenar grupos por valor central para reagrupar
            elementos_ordenados_centro = sorted(elementos, key=lambda e: (e[1]['Rango X'][0] + e[1]['Rango X'][1]) / 2)

            # Inicializar el grupo de centro con el primer elemento
            grupo_centro = [elementos_ordenados_centro[0]]

            for i in range(1, len(elementos_ordenados_centro)):
                actual_centro = elementos_ordenados_centro[i]
                anterior_centro = elementos_ordenados_centro[i - 1]

                # Comparar la distancia del centro de x con el anterior
                if abs((actual_centro[1]['Rango X'][0] + actual_centro[1]['Rango X'][1]) / 2 -
                       (anterior_centro[1]['Rango X'][0] + anterior_centro[1]['Rango X'][1]) / 2) < umbral_x:
                    grupo_centro.append(actual_centro)
                else:
                    # Crear un nuevo grupo para los centros y reiniciar
                    grupos[f'Grupo{grupo_id}'] = grupo_centro
                    grupo_centro = [actual_centro]
                    grupo_id += 1

            # Añadir el último grupo de centros
            if grupo_centro:
                grupos[f'Grupo{grupo_id}'] = grupo_centro

    # Construir el diccionario final con el número de elementos y textos ordenados
    resultado = {}
    for grupo, elementos in grupos.items():
        elementos_ordenados = sorted(elementos, key=lambda e: e[1]['Rango Y'][0])
        rangox = (min([el[1]['Rango X'][0] for el in elementos]), max([el[1]['Rango X'][1] for el in elementos]))
        rangoy = (min([el[1]['Rango Y'][0] for el in elementos]), max([el[1]['Rango Y'][1] for el in elementos]))
        colores_relleno_borde = {el[0]: (el[1]['Color Relleno'], el[1]['Color Borde']) for el in elementos}
        descripcion_general_colores, colores_excepcionales = describe_colores(colores_relleno_borde)
        textos = []
        for el in elementos_ordenados:
            if el[0] in colores_excepcionales.keys():
                texto = (el[0], f'\"{el[1]["Texto"]}\" ({colores_excepcionales[el[0]]})')
            else:
                texto = (el[0], f'\"{el[1]["Texto"]}\"')
            textos.append(texto)

        resultado[grupo] = {
            'Numero de elementos': len(elementos),
            'Textos': textos,
            'Posicion Relativa': obtener_posicion_relativa(rangox, rangoy, ancho_slide, alto_slide),
            'Descripcion Colores': descripcion_general_colores,
        }

    return resultado

def describe_traslape(dataframe_traslapes):
    imagenes = list(dataframe_traslapes['Imagen'])
    encimas = list(dataframe_traslapes['Encima'])
    traslapes = ['sobre' if el=='texto' else 'debajo de' for el in encimas]

    if len(imagenes) == 1:
        descripcion = f'{traslapes[0]} imagen {imagenes[0]}'

    elif len(imagenes) == 2:
        descripcion = f'{traslapes[0]} imagen {imagenes[0]} y {traslapes[1]} imagen {imagenes[1]}'
    else:
        descripcion = f'{traslapes[0]} imagen {imagenes[0]}'
        for ind in range(1,len(imagenes)-1):
            traslape, imagen = traslapes[ind], imagenes[ind]
            descripcion += (f', {traslape} imagen {imagen}')
        descripcion += (f'y {traslapes[ind]} imagen {imagenes[ind]}')
    return descripcion



def describe_picture_shape(n_slide, shape, n_shape,  ancho_slide, alto_slide):
    rangox = (shape.left/conversion, shape.left/conversion + shape.width/conversion)
    rangoy = (shape.top/conversion, shape.top/conversion + shape.height/conversion)
    nombre_archivo = f'slide{n_slide:03}_image{n_shape:02}.{shape.image.ext}'
    dict = {
        "Rango X": rangox,
        "Rango Y": rangoy,
        "Extension": f'{shape.image.ext}',
        "Orden": n_shape,
        "Posicion Relativa": obtener_posicion_relativa(rangox, rangoy, ancho_slide, alto_slide),
        "Nombre Archivo" : nombre_archivo,
        "Shape": shape
    }
    return dict

def describe_slide(n_slide,slide, ancho_slides, alto_slides):
    descripcion=''
    shapes = slide.shapes
    n_text_shape, n_picture_shape = 0,0
    text_shapes, picture_shapes = dict([]), dict([])
    for n_shape, shape in enumerate(shapes, start=1):
        if shape.shape_type == 13:
            n_picture_shape += 1
            picture_shapes[(f'imagen{n_picture_shape}.'
                        f'{shape.image.ext}')] = describe_picture_shape(n_slide,shape,n_shape,
                                                                        ancho_slides,alto_slides)
        if hasattr(shape, "text"):
            descripcion = describe_text_shape(shape, n_shape)
            ancho = descripcion['Rango X'][1]-descripcion['Rango X'][0]
            alto = descripcion['Rango Y'][1]-descripcion['Rango Y'][0]
            # Descartamos cuadros que enmarcan la diapositiva completa porque solo confunden
            if descripcion['Texto']=='' and ancho>0.7*ancho_slides and alto>0.7*alto_slides: continue
            n_text_shape +=1
            text_shapes[f'texto{n_text_shape}'] = descripcion
    df_traslapes, traslapes_mayoritarios = traslapes_imagen_texto(text_shapes, picture_shapes)
    df_traslapes_minoritarios = df_traslapes[~df_traslapes['Imagen'].isin(list(traslapes_mayoritarios.keys()))]
    grupos_texto = agrupar_text_shapes(text_shapes, ancho_slides, alto_slides)
    
    descripcion = ('A continuación se listan las componentes de texto e imagen de la slide\n\n')
    for __,grupo_texto in grupos_texto.items():
        n_elementos = grupo_texto['Numero de elementos']
        colores = grupo_texto['Descripcion Colores']
        posicion = grupo_texto['Posicion Relativa']
        textos = grupo_texto['Textos']
        if n_elementos>1:
            descripcion_grupo = (f'Hay un grupo de {n_elementos} elementos de texto {colores} cuya posicion '
                                 f'relativa es \"{posicion}\", y los textos de arriba a abajo son:\n')
            for texto_actual in textos:
                id_texto, texto = texto_actual
                # Revisar si hay información que dar sobre traslape con imagenes
                if id_texto in list(df_traslapes_minoritarios['Texto']):
                    info_traslapes = df_traslapes_minoritarios[df_traslapes_minoritarios['Texto']==id_texto]
                    texto_traslape = describe_traslape(info_traslapes)
                    if texto[-1] == ')': # Si ya existe un parentesis
                        descripcion_grupo += (f'{texto[:-1]}, ubicado {texto_traslape})\n')
                    else: # Si no hay todavia un paréntesis
                        descripcion_grupo += (f'{texto} (ubicado {texto_traslape})\n')
                else:
                    descripcion_grupo += (f'{texto}\n')
        
        else:
            descripcion_grupo = (f'Hay un texto {textos[0][1]} {colores} cuya posicion '
                                 f'relativa es \"{posicion}\"\n')
        descripcion += (f'\n{descripcion_grupo}\n')
        descripcion += '-'*40+'\n'

    if len(picture_shapes)>0:
        descripcion += (f'\nTambien están en la slide las siguientes imagenes:\n\n')
        for image_name, image_info in picture_shapes.items():
            posicion_relativa = image_info['Posicion Relativa']
            descripcion += (f'Imagen {image_name} cuya posición relativa es \"{posicion_relativa}\"\n')
    
    descripcion += '-'*40+'\n'

    return descripcion, picture_shapes


def guarda_slides(ppt_folder,output_folder,nombre_archivo,slides_prefix='slide'):

    pdf_path = f'{output_folder}/{nombre_archivo}.pdf'
    ppt_path = f'{ppt_folder}/{nombre_archivo}.pptx'

    subprocess.run(['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_folder, ppt_path])
    # Ejecuta soffice para convertir las diapositivas a imágenes PNG
    subprocess.run(['pdftoppm', '-png', '-rx', '96', '-ry', '96', pdf_path, os.path.join(output_folder, slides_prefix)])
    run_bash_command(f'rm -f {pdf_path}')





# Procesamiento de las presentaciones y conversión en lote
for nombre in nombres_archivos:
    nombre_sin_extension = nombre.split('.pptx')[0]
    print(' ')
    print(f'Extrayendo informacion de archivo {nombre_sin_extension}')
    dt = time.time() - ti
    print(f'Tiempo Transcurrido hasta ahora: {dt/60:.2f} minutos')
    
    run_bash_command(f'mkdir -p {path_salida}/{nombre_sin_extension}')
    run_bash_command(f'mkdir -p {path_slides}/{nombre_sin_extension}')
    
    guarda_slides(path_entrada,f'{path_slides}/{nombre_sin_extension}',nombre_sin_extension)

    prs = Presentation(f'{path_entrada}/{nombre}')
    ancho_slides = prs.slide_width/conversion
    alto_slides = prs.slide_height/conversion
    slides = prs.slides
    wmf_paths = []
    wmf_shapes = []
    print('Describiendo las Slides')
    for n_slide, slide in enumerate(slides, start=1):
        descripcion_slide, picture_shapes = describe_slide(n_slide, slide, ancho_slides, alto_slides)
        archivo_salida = open(f'{path_salida}/{nombre_sin_extension}/slide{n_slide:03}_content.txt', 'w')
        archivo_salida.write(descripcion_slide)
        archivo_salida.close()
        
        # Guardar y agrupar los WMF para convertir al final
        for nombre_imagen, info_imagen in picture_shapes.items():
            shape = info_imagen['Shape']
            nombre_archivo = info_imagen['Nombre Archivo']
            path_archivo = f'{path_salida}/{nombre_sin_extension}/{nombre_archivo}'
            extension = shape.image.ext
            if extension == 'wmf':
                guarda_wmf_shape_en_png(shape, path_archivo)
                wmf_paths.append(path_archivo.replace('.png', '.wmf'))
                wmf_shapes.append(shape)
            else:
                with open(path_archivo, 'wb') as archivo_imagen:
                    archivo_imagen.write(shape.image.blob)
                 # Revisar despues si efectivamente pasa para ambos la inversion
                if extension in ['jpeg','jpg']:
                    corregir_orientacion_imagen(path_archivo)


    # Convertir todos los WMF a PNG de una sola vez
    print('Conviertiendo archivos wmf a png')
    convertir_lote_wmf_a_png(wmf_paths)

    # Eliminar los archivos WMF temporales
    for wmf_path in wmf_paths:
        os.remove(wmf_path)


tf = time.time()
t_total = tf-ti
print(' ')
print(f'El script se demoró {t_total/60:.2f} minutos en extraer la información de {len(nombres_archivos)} archivos')

